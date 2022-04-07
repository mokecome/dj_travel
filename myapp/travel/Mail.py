from PIL import Image
import glob
import re
import shutil

def pptx_():
    # 處理圖片
    dest = shutil.copy('./myapp/travel/pic-eat.png', './myapp/travel/pic-e.png')
    pic = Image.open('./myapp/travel/pic-e.png')
    pic_copy =  pic.copy()
    width, height = pic_copy.size
    reset_size_img =pic_copy.resize((int(width*0.15), int(height*0.15)))
    reset_size_img.save('./myapp/travel/pic-e.png')
    
 

    # 檔名
    myfiles = glob.glob('./myapp/travel/*.png')  # 正則排除heart.png
    re_myfiles = []
    for i in myfiles:
        try:
            result = re.search('.+-.+', i)
            re_myfiles.append(result.group())
        except:
            pass
    z0 = list(map(lambda x: x.split('-')[0], re_myfiles))
    z1 = list(map(lambda x: x.split('-')[1][:-4], re_myfiles))
    dic_z = dict(zip(z0, z1))

    dest = shutil.copy('./myapp/travel/area-三亚.png', './myapp/travel/area.png')
    pic1 = Image.open('./myapp/travel/area.png')
    pic1_copy = pic1.copy()
    width, height = pic1_copy.size
    reset_size_img_1 =pic1_copy.resize((int(width *0.35), int(height*0.55)))
    reset_size_img_1.save('./myapp/travel/area.png')
    

    dest = shutil.copy('./myapp/travel/play-三五好友.png', './myapp/travel/play.png')
    pic2 =Image.open('./myapp/travel/play.png')
    pic2_copy = pic2.copy()
    width, height = pic2_copy.size
    reset_size_img_2 =pic2_copy.resize((int(width *0.75), int(height*0.75)))
    reset_size_img_2.save('./myapp/travel/play.png')
   

    from pptx import Presentation
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches
    from pptx.util import Cm, Pt

    prs = Presentation()
    """添加图片
    slide.shapes.add_picture(图片路径, 距离左边，距离顶端,宽度，高度)
    """
    blank_slide_layout = prs.slide_layouts[6]

    left = top = Cm(0.5)

    slide4 = prs.slides.add_slide(blank_slide_layout)
    pic4 = slide4.shapes.add_picture('./myapp/travel/pic-e.png', left, top)

    slide5 = prs.slides.add_slide(blank_slide_layout)
    pic5 = slide5.shapes.add_picture('./myapp/travel/area.png', left, top)

    slide6 = prs.slides.add_slide(blank_slide_layout)
    pic6 = slide6.shapes.add_picture('./myapp/travel/play.png', left, top)

    # 选择使用第二个模板
    bullet_slide_layout = prs.slide_layouts[1]
    # 使用模板生成一张幻灯片
    slide1 = prs.slides.add_slide(bullet_slide_layout)
    # 获取幻灯片的所有图形对象
    shapes = slide1.shapes

    # 设置标题
    title_shape = shapes.title
    # 获取第一个图片对象
    body_shape = shapes.placeholders[1]
    title_shape.text = '結論'

    tf = body_shape.text_frame
    try:
        t = int(dic_z['time'])
    except:
        t=5

    tf.text = '旅游高峰期在"{}"月，選擇避開高峰期的"{}"月和"{}"月的周末短途旅行也是不錯的選擇'.format(t, t - 1, t + 1)  # time-
    # 设定层级关系
    # paragraph.level=层级数
    # 0为最顶层
    # tf.level = 0

    p = tf.add_paragraph()
    try:
        if dic_z['pic'] == 'eat':
            eat = '美食'
    except:
        eat = '美食'

    p.text = '"{}"已成為旅游的代名詞'.format(eat)
    # p.level = 1

    p = tf.add_paragraph()
    try:
        dic_z['area']
    except:
        dic_z['area']='三亞'
    p.text = '個人認為性價比較高的旅游城市："{}"'.format(dic_z['area'])  # area-

    try:
        dic_z['play']
    except:
        dic_z['play']='三五好友'
    p = tf.add_paragraph()
    p.text = '"{}"一起旅游是最令人們喜歡的出游方式'.format(dic_z['play'])  # play-

    # p.level = 2

    """添加表格  
    """
    # import pandas as pd
    # df=pd.read_csv('data9.csv')

    # blank_slide_layout = prs.slide_layouts[6]
    # slide5 = prs.slides.add_slide(blank_slide_layout)
    # if isinstance(df,pd.Series):
    #     rows, cols = len(df.index),2
    # else:
    #     rows, cols = len(df.index),len(df.columns)

    # print()
    # left = top = Cm(5)
    # width = Cm(18)
    # height = Cm(3)

    # table = slide5.shapes.add_table(rows, cols, top, left, width, height).table

    # # table.columns[0].width = Cm(6)
    # # table.columns[1].width = Cm(8)
    # # table.rows[0].width = Cm(3)
    # #如果是
    # data = list(df.items())

    # #表頭
    # for row in range(rows):
    #     for col in range(cols):
    #         table.cell(row, col).text = str(data[row][col])

    prs.save('travel.pptx')





'''
receiver_email_=''
'''
import email, smtplib, ssl
from email import encoders
from email.mime.base import MIMEBase
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
def sendmail(mail):
    pptx_()
    subject = "歡迎註冊"
    body = "已完成註冊，感謝您"
    sender_email ="testdjangosendmail@gmail.com"
    receiver_email= mail
    password = "1qaz2wsx+-*/"
    
    # Create a multipart message and set headers
    message = MIMEMultipart()
    message["From"] = sender_email
    message["To"] = receiver_email
    message["Subject"] = subject
    message["Bcc"] = receiver_email  # Recommended for mass emails
    
    # Add body to email
    message.attach(MIMEText(body, "plain"))
    
    filename = "./myapp/travel/travel.pptx"  # In same directory as script

    # Open file in binary mode
    with open(filename, "rb") as attachment:
        # Add file as application/octet-stream
        # Email client can usually download this automatically as attachment
        part = MIMEBase("application", "octet-stream")
        part.set_payload(attachment.read())
    
    # Encode file in ASCII characters to send by email    
    encoders.encode_base64(part)
    
    # Add header as key/value pair to attachment part
    part.add_header(
        "Content-Disposition",
        f"attachment; filename= {filename}",
    )
    
    # Add attachment to message and convert message to string
    message.attach(part)
    text = message.as_string()
    
    # Log in to server using secure context and send email
    context = ssl.create_default_context()
    with smtplib.SMTP_SSL("smtp.gmail.com", 465, context=context) as server:
        server.login(sender_email, password)
        server.sendmail(sender_email, receiver_email, text)



pptx_()

