import cv2
import numpy as np
import glob
import re
import os
import shutil 
#處理圖片

pic= cv2.imread('pic-eat.png')
roi  = cv2.resize(pic,(650,500))
cv2.imwrite('pic-eat.png',roi)




#檔名
myfiles = glob.glob('*.png')  #正則排除heart.png
re_myfiles =[]
for i in myfiles:
    try: 
        result=re.search('.+-.+', i)
        re_myfiles.append(result.group())
    except:
        pass
z0=list(map(lambda x: x.split('-')[0],re_myfiles))    
z1=list(map(lambda x: x.split('-')[1][:-4],re_myfiles))
dic_z=dict(zip(z0,z1))



dest = shutil.copy('area-三亚.png','area.png') 
pic1= cv2.imread('area.png')
top_y=0
bottom_y=800
left_x=200
right_x=2000
pic1 = pic1[top_y:bottom_y, left_x:right_x]
roi1  = cv2.resize(pic1,(650,500))
cv2.imwrite('area.png',roi1)


dest = shutil.copy('play-三五好友.png','play.png') 
pic2= cv2.imread('play.png')
roi2  = cv2.resize(pic2,(650,500))
cv2.imwrite('play.png',roi2)




from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.util import Cm, Pt


prs = Presentation()
"""添加图片
slide.shapes.add_picture(图片路径, 距离左边，距离顶端,宽度，高度)
"""
blank_slide_layout = prs.slide_layouts[6]


left = top = Cm(0.5)
'''
slide3 = prs.slides.add_slide(blank_slide_layout)
pic1 = slide3.shapes.add_picture('time-05.png', left, top)
'''
slide4 = prs.slides.add_slide(blank_slide_layout)
pic4 = slide4.shapes.add_picture('pic-eat.png', left, top)

slide5 = prs.slides.add_slide(blank_slide_layout)
pic5 = slide5.shapes.add_picture('area.png', left, top)

slide6 = prs.slides.add_slide(blank_slide_layout)
pic6= slide6.shapes.add_picture('play.png', left, top)


# 选择使用第二个模板
bullet_slide_layout = prs.slide_layouts[1]
# 使用模板生成一张幻灯片
slide1 = prs.slides.add_slide(bullet_slide_layout)
# 获取幻灯片的所有图形对象
shapes = slide1.shapes

# 设置标题
title_shape = shapes.title
# 获取第一个图片对象
body_shape = shapes.placeholders[1]
title_shape.text = '結論'

tf = body_shape.text_frame

t=int(dic_z['time'])

tf.text = '旅游高峰期在"{}"月，選擇避開高峰期的"{}"月和"{}"月的周末短途旅行也是不錯的選擇'.format(t,t-1,t+1) #time-
# 设定层级关系
# paragraph.level=层级数
# 0为最顶层
# tf.level = 0

p = tf.add_paragraph()
if dic_z['pic']=='eat':
    eat='美食'
p.text = '"{}"已成為旅游的代名詞'.format(eat)
# p.level = 1

p = tf.add_paragraph()
p.text = '個人認為性價比較高的旅游城市："{}"'.format(dic_z['area']) #area-

p = tf.add_paragraph()
p.text = '"{}"一起旅游是最令人們喜歡的出游方式'.format(dic_z['play']) #play-

# p.level = 2


"""添加表格  
"""
#import pandas as pd
#df=pd.read_csv('data9.csv')


# blank_slide_layout = prs.slide_layouts[6]
# slide5 = prs.slides.add_slide(blank_slide_layout)
# if isinstance(df,pd.Series):
#     rows, cols = len(df.index),2
# else:
#     rows, cols = len(df.index),len(df.columns)

# print()
# left = top = Cm(5)
# width = Cm(18)
# height = Cm(3)

# table = slide5.shapes.add_table(rows, cols, top, left, width, height).table

# # table.columns[0].width = Cm(6)
# # table.columns[1].width = Cm(8)
# # table.rows[0].width = Cm(3)
# #如果是
# data = list(df.items())

# #表頭
# for row in range(rows):
#     for col in range(cols):
#         table.cell(row, col).text = str(data[row][col])

prs.save('travel.pptx')

















